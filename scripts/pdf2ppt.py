# -*- coding: utf-8 -*-
import sys
import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
import tempfile

def convert_pdf_to_ppt(pdf_path, ppt_path):
    try:
        if not os.path.exists(pdf_path):
            print(f"Error: PDF file not found at {pdf_path}", file=sys.stderr)
            sys.exit(1)
            
        doc = fitz.open(pdf_path)
        total_pages = len(doc)
        
        # 创建空白的 PPTX 数据流骨架
        prs = Presentation()
        # 清理默认留存的第一张母版幻灯片（如果 pptx 自带的话通常无需处理新建即为0）
        
        # 为了避免最终生成的 PPT 体积如同黑洞般膨胀（原版输出往往几十页叠加 1.4G），
        # 必须大幅降低采样系数，并切换至强有损压缩架构。zoom=1 对应常态化屏幕观感 (72 DPI)。
        # 若需要更清晰，最高不建议超过 1.5。
        zoom = 1.25
        mat = fitz.Matrix(zoom, zoom)
        
        for i in range(total_pages):
            page = doc.load_page(i)
            # 渲染出位图
            pix = page.get_pixmap(matrix=mat, alpha=False)
            
            # 必须使用 .jpg 来激活 PyMuPDF 藏在暗处的高强度有损剥离算法。绝对杜绝使用带 alpha 无损图层的 .png
            with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_img:
                temp_img_name = temp_img.name
            
            # 使用更极客的底层数据存储钩子并启用 75% 质量的大幅度视觉欺骗压缩
            img_data = pix.tobytes("jpeg", jpg_quality=75)
            with open(temp_img_name, "wb") as f:
                f.write(img_data)
            
            # 以原 PDF 第一页的物理尺寸决定这套 PPTX 的基础画布比例
            if i == 0:
                prs.slide_width = Pt(page.rect.width)
                prs.slide_height = Pt(page.rect.height)
                
            blank_slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 将渲染出的超清截图完全铺满于当前独立切片
            slide.shapes.add_picture(temp_img_name, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            # 及时释放磁盘清理残渣
            os.remove(temp_img_name)
            
            # 这一步极其关键：它复刻了之前 PDF 转 Word 的伪装信号，以兼容全栈 Node.js 的正则探测引擎雷达！
            print(f"[{i+1}/{total_pages}] Converting page {i+1}...")
            # 关键：由于 Python 默认带缓存输出，必须强制把系统流吐出去，否则 Node 会漏接
            sys.stdout.flush() 

        # 封口下线，压实打包为幻灯片总集归档
        prs.save(ppt_path)
        doc.close()
        print(f"Success: {ppt_path}")
        sys.exit(0)
        
    except Exception as e:
        print(f"Error: {str(e)}", file=sys.stderr)
        sys.exit(1)

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python3 pdf2ppt.py <input_pdf_path> <output_ppt_path>", file=sys.stderr)
        sys.exit(1)
        
    input_pdf = sys.argv[1]
    output_ppt = sys.argv[2]
    convert_pdf_to_ppt(input_pdf, output_ppt)
